# import collections.abc 
# import config 
# assert collections
# import tkinter as tk
# from pptx import Presentation
# from pptx.util import Inches, Pt 
# import openai 
# from io import BytesIO
# import requests 
# from openai import OpenAI
# import os


# # API Token
# # openai.api_key = config.API_KEY 
# openai.api_key = os.environ["OPENAI_API_KEY"]

# def slide_generator(text, prs):
    
#     # Image generation section
#     prompt = f"Summarise the following text into a DALL-E image generation prompt:\n {text}"
#     model_engine = "gpt-4"
#     dlp = openai.chat.completions.create(
#         model=model_engine, 
#         messages = [{"role":"user", "content":"I will ask you a question"},
#                     {"role": "assistant", "content":"Ok"},
#                     {"role":"user", "content": prompt}],
 
#         max_tokens = 250, 
#         n=1, 
#         stop=None,
#         temperature = 0.8
#     )
#     dalle_prompt = dlp.choices[0].message.content.split()
#     # dalle_prompt = dlp.choices[0].text 
#     client = OpenAI()

#     response = openai.images.generate(
#         model="dall-e-3",
    
#     # response = openai.Image.create(
#         prompt = " ".join(dalle_prompt) + " Style: digital art", 
#         n=1,
#         size="1024x1024"
#     )
#     image_url = response.data[0].url
    
#     # Text generation section
#     prompt_text = f"Create a bullet point text for a Powerpoint"\
#         f"slide from the following text:\n {text}"

#     ppt = openai.chat.completions.create(
#         model = model_engine, 
#         messages= [{"role":"user", "content":"I will ask you a question"},
#                     {"role": "assistant", "content":"Ok"},
#                     {"role":"user", "content":f"{prompt_text}"}],
#         max_tokens = 1024,
#         n=1, 
#         stop=None, 
#         temperature= 0.8
#     ) 
#     ppt_text = ppt.choices[0].message.content.strip()
#     prompt_title = f"Create a title for a Powerpoint" \
#         f"slide from the following text: \n {text}"
        
#     ppt = openai.chat.completions.create(
#         model = model_engine, 
#         messages= [{"role":"user", "content":"I will ask you a question"},
#                     {"role":"assistant", "content":"Ok"},
#                     {"role": "user", "content": f"{prompt_title}"}                
#         ]
#     )  
#     ppt_header = ppt.choices[0].message.content.strip()
    
#     # Add a new slide to the presentation
#     slide = prs.slides.add_slide(prs.slide_layouts[1])
#     response = requests.get(image_url) # http request to get the generated image
#     img_bytes = BytesIO(response.content)
#     slide.shapes.add_picture(img_bytes, Inches(1), Inches(1))
    
#     # Add text box
#     txBox = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(4), Inches(4))

#     tf = txBox.text_frame
#     tf.text = ppt_text 
#     title_shape = slide.shapes.title
#     title_shape.text = ppt_header     
    
# def get_slides(): 
#     text = text_field.get("1.0", "end-1c")
#     paragraphs = text.split("\n\n")
#     prs = Presentation()
#     width = Pt(1920)
#     height = Pt(1080) 
#     prs.slide_width = width 
#     prs.slide_height = height
#     for paragraph in paragraphs: 
#         slide_generator(paragraph, prs)
#     prs.save("my_presentation_2.pptx")


# app = tk.Tk() 
# app.title("Create PPT Slides") 
# app.geometry("800x600")

# # create text field
# text_field = tk.Text(app)
# text_field.pack(fill="both", expand=True)
# text_field.configure(wrap="word", font=("Arial", 12))
# text_field.focus_set()

# # Create the button to create slides
# create_button = tk.Button(app, text="Create Slides", command=get_slides)
# create_button.pack()



# app.mainloop()

import collections.abc
import config

assert collections
import tkinter as tk
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import openai
from io import BytesIO
import requests

# API Token
openai.api_key = config.API_KEY


def slide_generator(text, prs):
    prompt = f"Summarize the following text to a DALL-E image generation " \
             f"prompt: \n {text}"

    model_engine = "gpt-4"
    dlp = openai.chat.completions.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=250,
        n=1,
        stop=None,
        temperature=0.8
    )

    dalle_prompt = " ".join(dlp.choices[0].message.content.strip())

    response = openai.images.generate(
        model="dall-e-3",
        prompt=dalle_prompt + " Style: digital art",
        n=1,
        size="1024x1024"
    )
    image_url = response.data[0].url

    prompt = f"Create a bullet point text for a Powerpoint" \
             f"slide from the following text: \n {text}"
    ppt = openai.chat.completions.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_text = " ".join(ppt.choices[0].message.content.strip())

    prompt = f"Create a title for a Powerpoint" \
             f"slide from the following text: \n {text}"
    ppt = openai.chat.completions.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_header = " ".join(ppt.choices[0].message.content.strip())

    # # Add a new slide to the presentation
    # slide = prs.slides.add_slide(prs.slide_layouts[1])

    # response = requests.get(image_url)
    # img_bytes = BytesIO(response.content)
    # slide.shapes.add_picture(img_bytes, Inches(1), Inches(1))

    # # Add text box
    # txBox = slide.shapes.add_textbox(Inches(7), Inches(4),
    #                                  Inches(6), Inches(5))
    # tf = txBox.text_frame
    # tf.text = ppt_text

    # title_shape = slide.shapes.title
    # title_shape.text = ppt_header
    
    # Add a slide
    slide_layout = prs.slide_layouts[5]  # Using a blank layout
    slide = prs.slides.add_slide(slide_layout)

    response = requests.get(image_url)
    img_bytes = BytesIO(response.content)
    slide.shapes.add_picture(img_bytes, Inches(1), Inches(1))
    
    # Define dimensions for the right half of the slide
    slide_from_left = Inches(15.7)  # Half of the standard 10-inch width slide
    slide_width_half = Inches(10.5)
    slide_height = Inches(7.5)  # Standard height for content placement
    margin_top = Inches(1.5)  # Adjust as needed for aesthetics
    header_height = Inches(3)
    
    # Add a textbox for the header in the right half of the slide
    header_box = slide.shapes.add_textbox(slide_from_left, margin_top, slide_width_half, header_height) # left, top, width, height
    header_frame = header_box.text_frame
    header_frame.text = ppt_header 
    
    header_frame.word_wrap = True  # Ensure long text wraps to new line

    # Center the header text horizontally
    for paragraph in header_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER

    # Adjust the font size and boldness of the header, if necessary
    header_frame.paragraphs[0].font.size = Pt(32)  # Example font size
    header_frame.paragraphs[0].font.bold = True

    # Add a textbox for the main text below the header in the right half of the slide
    text_box = slide.shapes.add_textbox(slide_from_left, margin_top + header_height, slide_width_half, slide_height - Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = ppt_text  # Replace with your main text
    
    # # Optionally set auto size for the text frame to ensure text fits well
    text_frame.word_wrap = True  # Ensure long text wraps to new line
    
    # Center the main text horizontally in the text frame
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER

    # Save the presentation
    prs.save('ReplaceMe_Chapter1.pptx')


def get_slides():
    text = text_field.get("1.0", "end-1c")
    paragraphs = text.split("\n\n")
    prs = Presentation()
    width = Pt(1920)
    height = Pt(1080)
    prs.slide_width = width
    prs.slide_height = height
    for paragraph in paragraphs:
        slide_generator(paragraph, prs)

    prs.save("jellyfish2.pptx")


app = tk.Tk()
app.title("Crate PPT Slides")
app.geometry("800x600")

# Create text field
text_field = tk.Text(app)
text_field.pack(fill="both", expand=True)
text_field.configure(wrap="word", font=("Arial", 12))
text_field.focus_set()

# Create the button to create slides
create_button = tk.Button(app, text="Create Slides", command=get_slides)
create_button.pack()

app.mainloop()

# Findings from Single LoRA Image Generation Results 
# - In general, generated image results are better (more closely follow the LoRA) when the specific LoRA tag is included in the text prompt. 
#     - Hence, the use of LoRAs WITH the respective LoRA tags is explored
#     - LoRA tags can be found on the respective LoRA model download page 
# - if text is present in the generated image, it is highly probable that the text does not make sense.

# Findings from Multiple LoRAs Image Generation Results 
# Results from merging multiple LoRAs tend to be beautiful, aesthetically-pleasing for a landscape text prompt.
# This could be due to
# - a wider variety of stylistic renditions can be applied on landscape subject matter (better than single object form)
# - fortunate suitability of num_inference_steps and guidance scale
#   - optimal num_inference_steps and guidance scale can be determined empirically (trial and error)